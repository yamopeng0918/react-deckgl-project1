import csv
import json
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation

from scripts.create_decision_tree_results_slide import _comparison_sentence, create_results_slide


def _metrics(accuracy, recalls, parameters):
    return {
        "accuracy": accuracy,
        "labels": list(range(8)),
        "support": {str(label): value for label, value in enumerate([1, 289, 1278, 910, 538, 19, 4, 0])},
        "recall": {str(label): recall for label, recall in enumerate(recalls)},
        "periods": {"train": [1995, 2023], "test": [2024, 2026]},
        "train_rows": 13617,
        "test_rows": 3039,
        "selected_parameters": parameters,
    }


class DecisionTreeResultsSlideTest(unittest.TestCase):
    def setUp(self):
        self.directory = tempfile.TemporaryDirectory()
        self.root = Path(self.directory.name)
        self.decision_tree = _metrics(
            0.2833168806,
            [0.0, 0.2629757785, 0.2527386541, 0.2857142857, 0.3550185874, 0.4736842105, 0.5, None],
            {"max_depth": 12, "min_samples_leaf": 1},
        )
        self.random_forest = _metrics(
            0.4389601843,
            [0.0, 0.1937716263, 0.4319248826, 0.3978021978, 0.6486988848, 0.6842105263, 0.5, None],
            {"n_estimators": 200, "max_depth": 12, "min_samples_leaf": 1},
        )
        self.decision_tree_path = self.root / "decision-tree.json"
        self.random_forest_path = self.root / "random-forest.json"
        self.comparison_path = self.root / "comparison.csv"
        self.output = self.root / "review.pptx"
        self.decision_tree_path.write_text(json.dumps(self.decision_tree), encoding="utf-8")
        self.random_forest_path.write_text(json.dumps(self.random_forest), encoding="utf-8")
        self._write_comparison()

    def tearDown(self):
        self.directory.cleanup()

    def _write_comparison(self, random_forest_accuracy=0.4389601843, random_forest_macro_recall=None):
        fields = ["model", "accuracy", "macro_recall"]
        for label in range(8):
            fields.extend([f"recall_{label}", f"support_{label}"])
        with self.comparison_path.open("w", encoding="utf-8", newline="") as stream:
            writer = csv.DictWriter(stream, fieldnames=fields)
            writer.writeheader()
            for model, metrics, accuracy in (
                ("decision_tree", self.decision_tree, self.decision_tree["accuracy"]),
                ("random_forest", self.random_forest, random_forest_accuracy),
            ):
                row = {
                    "model": model,
                    "accuracy": accuracy,
                    "macro_recall": sum(value for value in metrics["recall"].values() if value is not None) / 7,
                }
                if model == "random_forest" and random_forest_macro_recall is not None:
                    row["macro_recall"] = random_forest_macro_recall
                for label in range(8):
                    row[f"recall_{label}"] = "" if metrics["recall"][str(label)] is None else metrics["recall"][str(label)]
                    row[f"support_{label}"] = metrics["support"][str(label)]
                writer.writerow(row)

    def test_creates_one_widescreen_two_model_slide(self):
        create_results_slide(
            self.decision_tree_path,
            self.random_forest_path,
            self.comparison_path,
            self.output,
        )

        with zipfile.ZipFile(self.output) as package:
            self.assertIsNone(package.testzip())
            self.assertEqual(
                len([name for name in package.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]),
                1,
            )
        deck = Presentation(self.output)
        self.assertEqual(len(deck.slides), 1)
        self.assertAlmostEqual(deck.slide_width / deck.slide_height, 16 / 9, places=2)
        text = "\n".join(shape.text for shape in deck.slides[0].shapes if hasattr(shape, "text"))
        for expected in (
            "決策樹",
            "隨機森林",
            "Accuracy",
            "Macro Recall",
            "稀有類別",
            "28.33%",
            "30.43%",
            "43.90%",
            "40.81%",
            "共同評估｜訓練 1995–2023（13,617 筆）｜測試 2024–2026（3,039 筆）",
            "選定參數｜max_depth=12 · min_samples_leaf=1",
            "選定參數｜n_estimators=200 · max_depth=12 · min_samples_leaf=1",
            "比較結論｜隨機森林在時間外推測試的 Macro Recall 領先；Accuracy 亦由隨機森林領先。",
        ):
            self.assertIn(expected, text)
        for expected in (
            "0.0%（1）", "26.3%（289）", "25.3%（1,278）", "28.6%（910）",
            "35.5%（538）", "47.4%（19）", "50.0%（4）", "N/A（0）",
            "19.4%（289）", "43.2%（1,278）", "39.8%（910）", "64.9%（538）",
            "68.4%（19）",
        ):
            self.assertIn(expected, text)
        self.assertGreaterEqual(text.count("0.0%（1）"), 2)
        self.assertGreaterEqual(text.count("50.0%（4）"), 2)
        self.assertGreaterEqual(text.count("N/A（0）"), 2)

    def test_rejects_comparison_values_that_disagree_with_metrics(self):
        self._write_comparison(random_forest_accuracy=0.9)

        with self.assertRaisesRegex(ValueError, "comparison.*random_forest.*accuracy"):
            create_results_slide(
                self.decision_tree_path,
                self.random_forest_path,
                self.comparison_path,
                self.output,
            )

    def test_rejects_comparison_macro_recall_that_disagrees_with_metrics(self):
        self._write_comparison(random_forest_macro_recall=0.9)

        with self.assertRaisesRegex(ValueError, "comparison.*random_forest.*macro_recall"):
            create_results_slide(
                self.decision_tree_path,
                self.random_forest_path,
                self.comparison_path,
                self.output,
            )

    def test_comparison_sentence_uses_accuracy_to_break_macro_recall_tie(self):
        comparison = {
            "decision_tree": {"macro_recall": "0.4", "accuracy": "0.51"},
            "random_forest": {"macro_recall": "0.4", "accuracy": "0.62"},
        }

        self.assertEqual(
            _comparison_sentence(comparison),
            "比較結論｜測試 Macro Recall 平手；隨機森林在 Accuracy 領先。",
        )

    def test_comparison_sentence_says_neither_leads_when_both_metrics_tie(self):
        comparison = {
            "decision_tree": {"macro_recall": "0.4", "accuracy": "0.51"},
            "random_forest": {"macro_recall": "0.4", "accuracy": "0.51"},
        }

        self.assertEqual(
            _comparison_sentence(comparison),
            "比較結論｜測試 Macro Recall 與 Accuracy 皆平手，兩者沒有明確領先。",
        )


if __name__ == "__main__":
    unittest.main()
