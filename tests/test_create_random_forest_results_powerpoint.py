import json
import math
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from scripts.create_random_forest_results_powerpoint import (
    create_random_forest_results_powerpoint,
    load_and_validate_metrics,
)


def valid_metrics():
    support = [1, 289, 1278, 910, 538, 19, 4, 0]
    matrix = [[0] * 8 for _ in range(8)]
    for label, count in enumerate(support):
        matrix[label][label] = count
    return {
        "accuracy": 0.4389601843,
        "labels": list(range(8)),
        "support": {str(i): value for i, value in enumerate(support)},
        "recall": {str(i): value for i, value in enumerate([0.0, 0.1938, 0.4319, 0.3978, 0.6487, 0.6842, 0.5, None])},
        "confusion_matrix": matrix,
        "periods": {"train": [1995, 2023], "test": [2024, 2026]},
        "train_rows": 13617,
        "test_rows": 3039,
        "selected_parameters": {
            "n_estimators": 200,
            "max_depth": 12,
            "min_samples_leaf": 1,
            "max_features": "sqrt",
            "class_weight": "balanced_subsample",
            "random_state": 42,
            "validation_start_year": 2021,
            "validation_end_year": 2023,
            "validation_macro_recall": 0.4206,
            "validation_accuracy": 0.4583,
        },
    }


class RandomForestResultsPowerPointTest(unittest.TestCase):
    def setUp(self):
        self.directory = tempfile.TemporaryDirectory()
        self.root = Path(self.directory.name)
        self.metrics_path = self.root / "metrics.json"
        self.matrix_path = self.root / "matrix.png"
        self.output_path = self.root / "results.pptx"
        self.metrics_path.write_text(json.dumps(valid_metrics()), encoding="utf-8")
        Image.new("RGB", (1200, 900), "white").save(self.matrix_path)

    def tearDown(self):
        self.directory.cleanup()

    def test_creates_valid_two_slide_widescreen_deck(self):
        create_random_forest_results_powerpoint(
            self.metrics_path, self.matrix_path, self.output_path
        )

        with zipfile.ZipFile(self.output_path) as package:
            self.assertIsNone(package.testzip())
        deck = Presentation(self.output_path)
        self.assertEqual(len(deck.slides), 2)
        self.assertAlmostEqual(deck.slide_width / deck.slide_height, 16 / 9, places=2)
        text = "\n".join(
            shape.text
            for slide in deck.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
        for expected in (
            "隨機森林最大震度分類結果",
            "Accuracy",
            "Macro Recall",
            "1995–2023",
            "2024–2026",
            "13,617",
            "3,039",
            "n_estimators=200",
            "max_features=sqrt",
            "強度 7：N/A（support 0）",
            "稀有類別",
            "分類，不是地震預測",
            "混淆矩陣",
            "實際類別（縱軸）",
            "預測類別（橫軸）",
            "2–4 級相鄰混淆",
            "0、6、7 級",
        ):
            self.assertIn(expected, text)

        pictures = [
            shape for shape in deck.slides[1].shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertEqual(len(pictures), 1)
        picture = pictures[0]
        self.assertGreater(picture.width, deck.slide_width * 0.45)
        self.assertGreater(picture.height, deck.slide_height * 0.55)
        self.assertGreaterEqual(picture.left, 0)
        self.assertGreaterEqual(picture.top, 0)
        self.assertLessEqual(picture.left + picture.width, deck.slide_width)
        self.assertLessEqual(picture.top + picture.height, deck.slide_height)

    def test_assigns_required_content_to_the_correct_slide(self):
        create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)
        deck = Presentation(self.output_path)
        texts = ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in deck.slides]
        self.assertIn("Accuracy", texts[0])
        self.assertIn("強度 7：N/A（support 0）", texts[0])
        self.assertNotIn("2–4 級相鄰混淆", texts[0])
        self.assertIn("2–4 級相鄰混淆", texts[1])
        self.assertIn("實際類別（縱軸）", texts[1])
        self.assertNotIn("n_estimators=200", texts[1])
        for key, value in valid_metrics()["selected_parameters"].items():
            self.assertIn(f"{key}={value}", texts[0])

    def test_uses_jhenghei_and_established_palette(self):
        create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)
        deck = Presentation(self.output_path)
        fonts = {
            run.font.name
            for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text_frame")
            for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text
        }
        self.assertEqual(fonts, {"Microsoft JhengHei"})
        expected = {"132338", "1F807D", "DA9335", "F7F5EF"}
        with zipfile.ZipFile(self.output_path) as package:
            xml = b"".join(package.read(name) for name in package.namelist() if name.endswith(".xml"))
        for color in expected:
            self.assertIn(color.encode(), xml)

    def test_embeds_exact_uncropped_aspect_preserved_source_png(self):
        create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)
        deck = Presentation(self.output_path)
        picture = next(shape for shape in deck.slides[1].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
        self.assertEqual(picture.image.blob, self.matrix_path.read_bytes())
        self.assertEqual((picture.crop_left, picture.crop_right, picture.crop_top, picture.crop_bottom), (0, 0, 0, 0))
        self.assertAlmostEqual(picture.width / picture.height, 1200 / 900, places=4)

    def _write_metrics(self, metrics):
        self.metrics_path.write_text(json.dumps(metrics), encoding="utf-8")

    def test_rejects_labels_other_than_exactly_zero_through_seven(self):
        metrics = valid_metrics()
        metrics["labels"] = list(range(7))
        self._write_metrics(metrics)

        with self.assertRaisesRegex(ValueError, "labels.*0.*7"):
            load_and_validate_metrics(self.metrics_path)

    def test_rejects_confusion_matrix_that_is_not_eight_by_eight(self):
        metrics = valid_metrics()
        metrics["confusion_matrix"] = [[0] * 8 for _ in range(7)]
        self._write_metrics(metrics)

        with self.assertRaisesRegex(ValueError, "confusion matrix.*8.*8"):
            load_and_validate_metrics(self.metrics_path)

    def test_rejects_matrix_row_total_that_disagrees_with_support(self):
        metrics = valid_metrics()
        metrics["confusion_matrix"][2][2] -= 1
        self._write_metrics(metrics)

        with self.assertRaisesRegex(ValueError, "row 2.*support"):
            load_and_validate_metrics(self.metrics_path)

    def test_rejects_missing_png_before_altering_output(self):
        self.output_path.write_bytes(b"existing")
        self.matrix_path.unlink()

        with self.assertRaisesRegex(FileNotFoundError, "confusion matrix PNG"):
            create_random_forest_results_powerpoint(
                self.metrics_path, self.matrix_path, self.output_path
            )

        self.assertEqual(self.output_path.read_bytes(), b"existing")

    def test_save_failure_preserves_existing_output_and_removes_temporary_file(self):
        self.output_path.write_bytes(b"existing deck")

        class FailingDeck:
            def save(self, _path):
                raise OSError("simulated save failure")

        with patch(
            "scripts.create_random_forest_results_powerpoint._build_deck",
            return_value=FailingDeck(),
        ):
            with self.assertRaisesRegex(OSError, "simulated save failure"):
                create_random_forest_results_powerpoint(
                    self.metrics_path, self.matrix_path, self.output_path
                )

        self.assertEqual(self.output_path.read_bytes(), b"existing deck")
        self.assertEqual(list(self.root.glob(".results-*.pptx")), [])

    def test_rejects_invalid_recall_values_and_key_sets(self):
        for name, mutate in (
            ("keys", lambda m: m["recall"].pop("7")),
            ("bool", lambda m: m["recall"].__setitem__("1", True)),
            ("negative", lambda m: m["recall"].__setitem__("1", -0.1)),
            ("over-one", lambda m: m["recall"].__setitem__("1", 1.1)),
            ("nonfinite", lambda m: m["recall"].__setitem__("1", math.inf)),
        ):
            with self.subTest(name=name):
                metrics = valid_metrics(); mutate(metrics); self._write_metrics(metrics)
                with self.assertRaisesRegex(ValueError, "recall"):
                    load_and_validate_metrics(self.metrics_path)

    def test_rejects_invalid_support_values_and_key_sets(self):
        for name, mutate in (
            ("keys", lambda m: m["support"].pop("7")),
            ("bool", lambda m: m["support"].__setitem__("1", True)),
            ("negative", lambda m: m["support"].__setitem__("1", -1)),
            ("fractional", lambda m: m["support"].__setitem__("1", 289.5)),
            ("nonfinite", lambda m: m["support"].__setitem__("1", math.inf)),
        ):
            with self.subTest(name=name):
                metrics = valid_metrics(); mutate(metrics); self._write_metrics(metrics)
                with self.assertRaisesRegex(ValueError, "support"):
                    load_and_validate_metrics(self.metrics_path)

    def test_rejects_invalid_matrix_cell_values(self):
        for name, value in (("bool", True), ("negative", -1), ("fractional", 0.5), ("nonfinite", math.nan)):
            with self.subTest(name=name):
                metrics = valid_metrics(); metrics["confusion_matrix"][0][0] = value; self._write_metrics(metrics)
                with self.assertRaisesRegex(ValueError, "confusion matrix"):
                    load_and_validate_metrics(self.metrics_path)

    def test_rejects_invalid_accuracy_periods_rows_and_parameters(self):
        cases = (
            ("accuracy", lambda m: m.__setitem__("accuracy", True), "accuracy"),
            ("accuracy-range", lambda m: m.__setitem__("accuracy", math.inf), "accuracy"),
            ("period-shape", lambda m: m["periods"].__setitem__("train", [1995]), "periods"),
            ("period-order", lambda m: m["periods"].__setitem__("test", [2026, 2024]), "periods"),
            ("period-bool", lambda m: m["periods"].__setitem__("test", [True, 2026]), "periods"),
            ("train-rows", lambda m: m.__setitem__("train_rows", 0), "train_rows"),
            ("test-rows", lambda m: m.__setitem__("test_rows", 3.5), "test_rows"),
            ("parameters-dict", lambda m: m.__setitem__("selected_parameters", []), "selected_parameters"),
            ("parameter-key", lambda m: m["selected_parameters"].pop("max_features"), "selected_parameters"),
        )
        for name, mutate, message in cases:
            with self.subTest(name=name):
                metrics = valid_metrics(); mutate(metrics); self._write_metrics(metrics)
                with self.assertRaisesRegex(ValueError, message):
                    load_and_validate_metrics(self.metrics_path)

    def test_rejects_unreadable_and_non_png_images(self):
        for name, payload in (("unreadable", b"not an image"), ("jpeg", None)):
            with self.subTest(name=name):
                if payload is None:
                    Image.new("RGB", (10, 10)).save(self.matrix_path, format="JPEG")
                else:
                    self.matrix_path.write_bytes(payload)
                with self.assertRaisesRegex(ValueError, "PNG"):
                    create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)

    def test_corrupt_generated_package_preserves_output_and_cleans_temp(self):
        self.output_path.write_bytes(b"existing")
        class CorruptDeck:
            def save(self, path): Path(path).write_bytes(b"not a PowerPoint package")
        with patch("scripts.create_random_forest_results_powerpoint._build_deck", return_value=CorruptDeck()):
            with self.assertRaisesRegex(ValueError, "corrupt"):
                create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)
        self.assertEqual(self.output_path.read_bytes(), b"existing")
        self.assertEqual(list(self.root.glob(".results-*.pptx")), [])

    def test_atomic_replace_failure_preserves_output_and_cleans_temp(self):
        self.output_path.write_bytes(b"existing")
        with patch("scripts.create_random_forest_results_powerpoint.os.replace", side_effect=OSError("replace failed")):
            with self.assertRaisesRegex(OSError, "replace failed"):
                create_random_forest_results_powerpoint(self.metrics_path, self.matrix_path, self.output_path)
        self.assertEqual(self.output_path.read_bytes(), b"existing")
        self.assertEqual(list(self.root.glob(".results-*.pptx")), [])


if __name__ == "__main__":
    unittest.main()
