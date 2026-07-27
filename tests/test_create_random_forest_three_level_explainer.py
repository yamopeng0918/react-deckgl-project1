from collections import Counter
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch
from xml.etree import ElementTree

import joblib
from pptx import Presentation
from sklearn.ensemble import RandomForestClassifier

from scripts.create_random_forest_three_level_explainer import (
    FEATURE_LABELS,
    FEATURE_NAMES,
    build_deck,
    create_random_forest_three_level_explainer,
    extract_three_levels,
    select_representative_tree,
)


def fitted_model():
    features = []
    labels = []
    for repetition in range(6):
        for intensity in range(8):
            features.append(
                [
                    3.0 + intensity + repetition / 100,
                    5.0 + (intensity % 3) + repetition / 10,
                    120.0 + intensity / 10,
                    22.0 + repetition / 10,
                    float((intensity % 12) + 1),
                    float((repetition * 3 + intensity) % 24),
                ]
            )
            labels.append(intensity)

    model = RandomForestClassifier(
        n_estimators=200,
        max_depth=4,
        random_state=42,
    )
    model.fit(features, labels)
    return model


class RepresentativeTreeExtractionTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.model = fitted_model()

    def test_selection_is_deterministic(self):
        index_a, tree_a = select_representative_tree(self.model)
        index_b, tree_b = select_representative_tree(self.model)

        self.assertEqual(index_a, index_b)
        self.assertIs(tree_a, tree_b)

    def test_feature_labels_match_the_presentation_contract(self):
        self.assertEqual(
            FEATURE_LABELS,
            {
                "magnitude": "\u9707\u7d1a",
                "depth_km": "\u6df1\u5ea6",
                "longitude": "\u7d93\u5ea6",
                "latitude": "\u7def\u5ea6",
                "month": "\u6708\u4efd",
                "hour": "\u6642\u523b",
            },
        )

    def test_extraction_matches_estimator_values_at_first_three_levels(self):
        _, tree = select_representative_tree(self.model)

        nodes = extract_three_levels(tree, FEATURE_NAMES, list(self.model.classes_))

        self.assertLessEqual(len(nodes), 7)
        self.assertEqual(nodes[0]["depth"], 0)
        self.assertTrue(all(node["depth"] <= 2 for node in nodes))
        self.assertEqual(
            [node["node_id"] for node in nodes],
            sorted(node["node_id"] for node in nodes),
        )
        self.assertTrue(
            any(node["depth"] == 2 and not node["is_leaf"] for node in nodes)
        )

        for node in nodes:
            tree_node = tree.tree_
            expected_counts = tuple(
                float(value) for value in tree_node.value[node["node_id"]][0]
            )
            expected_dominant_class = int(
                self.model.classes_[
                    max(range(len(expected_counts)), key=expected_counts.__getitem__)
                ]
            )
            self.assertEqual(
                node["samples"], int(tree_node.n_node_samples[node["node_id"]])
            )
            self.assertEqual(node["dominant_class"], expected_dominant_class)
            self.assertEqual(node["class_counts"], expected_counts)

            if not node["is_leaf"]:
                self.assertEqual(
                    node["feature_index"], int(tree_node.feature[node["node_id"]])
                )
                self.assertEqual(
                    node["threshold"], float(tree_node.threshold[node["node_id"]])
                )

        positions = {node["node_id"]: index for index, node in enumerate(nodes)}
        for node in nodes:
            if node["parent_id"] is not None:
                self.assertLess(positions[node["parent_id"]], positions[node["node_id"]])
        for parent_id in {
            node["parent_id"] for node in nodes if node["parent_id"] is not None
        }:
            children = [node for node in nodes if node["parent_id"] == parent_id]
            if len(children) == 2:
                self.assertEqual([child["branch"] for child in children], ["yes", "no"])


class RepresentativeTreeValidationTest(unittest.TestCase):
    def test_rejects_model_with_wrong_feature_count(self):
        model = RandomForestClassifier(n_estimators=1, random_state=42)
        model.fit([[0, 1, 2, 3, 4], [1, 2, 3, 4, 5]], [0, 1])

        with self.assertRaisesRegex(ValueError, "six input features"):
            select_representative_tree(model)


class RandomForestThreeLevelExplainerPowerPointTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.model = fitted_model()

    def setUp(self):
        self.directory = tempfile.TemporaryDirectory()
        self.root = Path(self.directory.name)
        self.model_path = self.root / "model.joblib"
        self.output_path = self.root / "explainer.pptx"
        joblib.dump(self.model, self.model_path)

    def tearDown(self):
        self.directory.cleanup()

    @staticmethod
    def _slide_text(slide):
        return "\n".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )

    @staticmethod
    def _expected_node_texts(nodes):
        expected = []
        for node in nodes:
            if node["is_leaf"]:
                expected.append(
                    f"此分支預測震度 {node['dominant_class']}\n"
                    f"樣本 {node['samples']:,}"
                )
            else:
                expected.append(
                    f"{FEATURE_LABELS[node['feature_name']]} ≤ "
                    f"{node['threshold']:.2f}？\n"
                    f"樣本 {node['samples']:,}｜目前偏向震度 "
                    f"{node['dominant_class']}"
                )
        return expected

    def test_build_deck_displays_every_extracted_node_value(self):
        tree_index, tree = select_representative_tree(self.model)
        nodes = extract_three_levels(
            tree,
            FEATURE_NAMES,
            list(self.model.classes_),
        )

        deck = build_deck(nodes, tree_index, len(self.model.estimators_))

        self.assertEqual(len(deck.slides), 1)
        expected = self._expected_node_texts(nodes)
        node_shapes = [
            shape
            for shape in deck.slides[0].shapes
            if hasattr(shape, "text")
            and (
                shape.text.startswith("此分支預測震度 ")
                or "｜目前偏向震度 " in shape.text
            )
        ]
        self.assertLessEqual(len(node_shapes), 7)
        self.assertEqual(len(node_shapes), len(nodes))
        self.assertEqual(
            Counter(shape.text for shape in node_shapes),
            Counter(expected),
        )

    def test_creates_valid_one_slide_widescreen_native_shape_deck(self):
        result = create_random_forest_three_level_explainer(
            self.model_path,
            self.output_path,
        )

        self.assertEqual(result, self.output_path)
        with zipfile.ZipFile(self.output_path) as package:
            self.assertIsNone(package.testzip())
        deck = Presentation(self.output_path)
        self.assertEqual(len(deck.slides), 1)
        self.assertAlmostEqual(
            deck.slide_width / deck.slide_height,
            16 / 9,
            places=2,
        )
        slide = deck.slides[0]
        slide_text = self._slide_text(slide)
        for expected in (
            "隨機森林如何判斷最大震度？",
            "是（≤）",
            "否（>）",
            "200 棵樹中的一棵",
            "投票",
            "不是地震預測",
        ):
            self.assertIn(expected, slide_text)

        for shape in slide.shapes:
            with self.subTest(shape=shape.name):
                self.assertGreaterEqual(shape.left, 0)
                self.assertGreaterEqual(shape.top, 0)
                self.assertLessEqual(shape.left + shape.width, deck.slide_width)
                self.assertLessEqual(shape.top + shape.height, deck.slide_height)

    def test_serializes_jhenghei_as_the_east_asian_font_for_every_text_run(self):
        create_random_forest_three_level_explainer(
            self.model_path,
            self.output_path,
        )

        with zipfile.ZipFile(self.output_path) as package:
            slide_xml = ElementTree.fromstring(
                package.read("ppt/slides/slide1.xml")
            )
        namespace = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
        }
        runs = slide_xml.findall(".//a:r", namespace)
        self.assertGreater(len(runs), 0)
        for run in runs:
            with self.subTest(text=run.findtext("a:t", default="", namespaces=namespace)):
                east_asian = run.find("a:rPr/a:ea", namespace)
                self.assertIsNotNone(east_asian)
                self.assertEqual(
                    east_asian.attrib["typeface"],
                    "Microsoft JhengHei",
                )

    def test_save_failure_preserves_existing_output_and_cleans_temporary_file(self):
        self.output_path.write_bytes(b"existing deck")

        class FailingDeck:
            def save(self, _path):
                raise OSError("simulated save failure")

        with patch(
            "scripts.create_random_forest_three_level_explainer.build_deck",
            return_value=FailingDeck(),
        ):
            with self.assertRaisesRegex(OSError, "simulated save failure"):
                create_random_forest_three_level_explainer(
                    self.model_path,
                    self.output_path,
                )

        self.assertEqual(self.output_path.read_bytes(), b"existing deck")
        self.assertEqual(list(self.root.glob(".explainer-*.pptx")), [])

    def test_invalid_generated_package_preserves_output_and_cleans_temporary_file(self):
        self.output_path.write_bytes(b"existing deck")

        class CorruptDeck:
            def save(self, path):
                Path(path).write_bytes(b"not a PowerPoint package")

        with patch(
            "scripts.create_random_forest_three_level_explainer.build_deck",
            return_value=CorruptDeck(),
        ):
            with self.assertRaisesRegex(ValueError, "corrupt"):
                create_random_forest_three_level_explainer(
                    self.model_path,
                    self.output_path,
                )

        self.assertEqual(self.output_path.read_bytes(), b"existing deck")
        self.assertEqual(list(self.root.glob(".explainer-*.pptx")), [])


if __name__ == "__main__":
    unittest.main()
