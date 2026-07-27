import math
import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch
from xml.etree import ElementTree

import joblib
from pptx import Presentation
from pptx.oxml.ns import qn

from scripts.create_random_forest_feature_importance_powerpoint import (
    FEATURE_LABELS,
    FEATURE_NAMES,
    build_deck,
    build_insights,
    create_feature_importance_powerpoint,
    load_feature_importances,
    rank_feature_importances,
)


CURRENT_IMPORTANCES = (
    ("magnitude", 0.3113),
    ("depth_km", 0.1596),
    ("longitude", 0.1406),
    ("latitude", 0.2180),
    ("month", 0.0753),
    ("hour", 0.0952),
)


class FeatureImportanceExtractionTest(unittest.TestCase):
    def setUp(self):
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary_directory.name)
        self.model_path = self.root / "model.joblib"

    def tearDown(self):
        self.temporary_directory.cleanup()

    def _write_model(self, **attributes):
        joblib.dump(SimpleNamespace(**attributes), self.model_path)

    def test_loads_the_six_feature_importances_in_model_feature_order(self):
        self._write_model(
            feature_importances_=[value for _, value in CURRENT_IMPORTANCES]
        )

        values = load_feature_importances(self.model_path)

        self.assertEqual(tuple(name for name, _ in values), FEATURE_NAMES)
        self.assertAlmostEqual(sum(value for _, value in values), 1.0)
        self.assertIsInstance(values, tuple)

    def test_rejects_a_missing_model_file(self):
        with self.assertRaisesRegex(FileNotFoundError, "model file not found"):
            load_feature_importances(self.model_path)

    def test_corrupt_model_error_identifies_path_and_preserves_cause(self):
        self.model_path.write_bytes(b"not a readable joblib model")

        with self.assertRaises(ValueError) as raised:
            load_feature_importances(self.model_path)

        self.assertIn(str(self.model_path), str(raised.exception))
        self.assertIsNotNone(raised.exception.__cause__)

    def test_rejects_a_model_without_feature_importances(self):
        self._write_model()

        with self.assertRaisesRegex(ValueError, "feature_importances_"):
            load_feature_importances(self.model_path)

    def test_rejects_any_feature_count_other_than_six(self):
        self._write_model(feature_importances_=[0.2] * 5)

        with self.assertRaisesRegex(ValueError, "six feature importances"):
            load_feature_importances(self.model_path)

    def test_rejects_negative_or_non_finite_importances(self):
        for values in (
            [-0.01, 0.2, 0.2, 0.2, 0.2, 0.21],
            [math.nan, 0.2, 0.2, 0.2, 0.2, 0.2],
            [math.inf, 0.2, 0.2, 0.2, 0.2, 0.2],
        ):
            with self.subTest(values=values):
                self._write_model(feature_importances_=values)

                with self.assertRaisesRegex(ValueError, "finite non-negative"):
                    load_feature_importances(self.model_path)

    def test_rejects_importances_that_do_not_sum_to_one(self):
        self._write_model(feature_importances_=[0.1] * 6)

        with self.assertRaisesRegex(ValueError, "sum to one"):
            load_feature_importances(self.model_path)


class FeatureImportanceRankingTest(unittest.TestCase):
    def test_ranks_descending_and_preserves_feature_order_for_ties(self):
        values = (
            ("magnitude", 0.2),
            ("depth_km", 0.3),
            ("longitude", 0.3),
            ("latitude", 0.1),
            ("month", 0.05),
            ("hour", 0.05),
        )

        ranked = rank_feature_importances(values)

        self.assertEqual(
            ranked,
            (
                ("depth_km", 0.3),
                ("longitude", 0.3),
                ("magnitude", 0.2),
                ("latitude", 0.1),
                ("month", 0.05),
                ("hour", 0.05),
            ),
        )
        self.assertIsInstance(ranked, tuple)


class FeatureImportanceInsightTest(unittest.TestCase):
    def test_builds_current_model_insights_with_exact_percentages(self):
        insights = build_insights(CURRENT_IMPORTANCES)

        self.assertEqual(
            insights,
            (
                "規模是最重要的單一特徵（31.13%）",
                "經緯度合計 35.86%，顯示地理位置整體影響高於單一規模",
                "月份與時刻合計 17.05%，時間訊號存在，但不是主要依據",
            ),
        )
        self.assertIsInstance(insights, tuple)

    def test_uses_neutral_spatial_language_when_spatial_total_does_not_exceed_top_feature(self):
        insights = build_insights(
            (
                ("magnitude", 0.5),
                ("depth_km", 0.1),
                ("longitude", 0.1),
                ("latitude", 0.1),
                ("month", 0.1),
                ("hour", 0.1),
            )
        )

        self.assertNotIn("高於單一規模", insights[1])

    def test_uses_displayed_precision_for_spatial_comparison_boundaries(self):
        cases = (
            (0.30004, "30.00%", False),
            (0.30006, "30.01%", True),
        )
        for spatial_total, displayed_total, should_say_higher in cases:
            with self.subTest(spatial_total=spatial_total):
                insights = build_insights(
                    (
                        ("magnitude", 0.30001),
                        ("depth_km", 0.19999),
                        ("longitude", 0.15),
                        ("latitude", spatial_total - 0.15),
                        ("month", 0.1),
                        ("hour", 0.1),
                    )
                )

                self.assertIn(displayed_total, insights[1])
                self.assertEqual(
                    "高於單一規模" in insights[1],
                    should_say_higher,
                )

    def test_does_not_claim_a_temporal_signal_when_total_is_zero(self):
        insights = build_insights(
            (
                ("magnitude", 0.4),
                ("depth_km", 0.2),
                ("longitude", 0.2),
                ("latitude", 0.2),
                ("month", 0.0),
                ("hour", 0.0),
            )
        )

        self.assertIn("合計 0.00%", insights[2])
        self.assertNotIn("時間訊號存在", insights[2])

    def test_uses_displayed_precision_for_temporal_comparison_boundaries(self):
        cases = (
            (0.20004, "20.00%", True),
            (0.20006, "20.01%", False),
        )
        for temporal_total, displayed_total, should_say_not_primary in cases:
            with self.subTest(temporal_total=temporal_total):
                insights = build_insights(
                    (
                        ("magnitude", 0.20001),
                        ("depth_km", 0.29999),
                        ("longitude", 0.15),
                        ("latitude", 0.15),
                        ("month", 0.1),
                        ("hour", temporal_total - 0.1),
                    )
                )

                self.assertIn(displayed_total, insights[2])
                self.assertEqual(
                    "不是主要依據" in insights[2],
                    should_say_not_primary,
                )

    def test_uses_neutral_temporal_language_when_temporal_total_is_unusually_high(self):
        insights = build_insights(
            (
                ("magnitude", 0.25),
                ("depth_km", 0.1),
                ("longitude", 0.15),
                ("latitude", 0.1),
                ("month", 0.25),
                ("hour", 0.15),
            )
        )

        self.assertNotIn("不是主要依據", insights[2])


class FeatureImportancePowerPointTest(unittest.TestCase):
    def setUp(self):
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary_directory.name)
        self.model_path = self.root / "model.joblib"
        self.output_path = self.root / "feature-importance.pptx"
        joblib.dump(
            SimpleNamespace(
                feature_importances_=[
                    value for _, value in CURRENT_IMPORTANCES
                ]
            ),
            self.model_path,
        )

    def tearDown(self):
        self.temporary_directory.cleanup()

    @staticmethod
    def _slide_text(slide):
        return "\n".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )

    @staticmethod
    def _shape(slide, name):
        matches = [shape for shape in slide.shapes if shape.name == name]
        if len(matches) != 1:
            raise AssertionError(
                f"expected exactly one shape named {name!r}, got {len(matches)}"
            )
        return matches[0]

    def test_builds_one_slide_with_all_ranked_features_and_insights(self):
        ranked = rank_feature_importances(CURRENT_IMPORTANCES)
        insights = build_insights(CURRENT_IMPORTANCES)

        deck = build_deck(ranked, insights)

        self.assertEqual(len(deck.slides), 1)
        self.assertAlmostEqual(
            deck.slide_width / deck.slide_height,
            16 / 9,
            places=2,
        )
        slide_text = self._slide_text(deck.slides[0])
        for label in ("規模", "深度", "經度", "緯度", "月份", "時刻"):
            self.assertIn(label, slide_text)
        for insight in insights:
            self.assertIn(insight, slide_text)

    def test_bars_encode_descending_importances_from_a_shared_zero_baseline(self):
        ranked = rank_feature_importances(CURRENT_IMPORTANCES)
        deck = build_deck(ranked, build_insights(CURRENT_IMPORTANCES))
        slide = deck.slides[0]

        bars = [
            self._shape(slide, f"importance-bar-{feature_name}")
            for feature_name, _ in ranked
        ]
        labels = [
            self._shape(slide, f"importance-label-{feature_name}")
            for feature_name, _ in ranked
        ]
        self.assertEqual(len(bars), 6)
        self.assertEqual([shape.top for shape in bars], sorted(shape.top for shape in bars))
        self.assertEqual(
            [shape.top for shape in labels],
            sorted(shape.top for shape in labels),
        )
        self.assertEqual(len({shape.left for shape in bars}), 1)

        maximum_width = bars[0].width
        maximum_importance = ranked[0][1]
        for bar, (_, value) in zip(bars[1:], ranked[1:]):
            self.assertAlmostEqual(
                bar.width / maximum_width,
                value / maximum_importance,
                places=6,
            )

        for feature_name, value in ranked:
            value_shape = self._shape(
                slide,
                f"importance-value-{feature_name}",
            )
            self.assertEqual(value_shape.text, f"{value:.2%}")
            self.assertEqual(
                self._shape(
                    slide,
                    f"importance-label-{feature_name}",
                ).text,
                FEATURE_LABELS[feature_name],
            )

    def test_percentage_labels_follow_their_associated_bar_ends(self):
        ranked = rank_feature_importances(CURRENT_IMPORTANCES)
        deck = build_deck(ranked, build_insights(CURRENT_IMPORTANCES))
        slide = deck.slides[0]

        end_gaps = []
        for feature_name, _ in ranked:
            bar = self._shape(slide, f"importance-bar-{feature_name}")
            value_label = self._shape(
                slide,
                f"importance-value-{feature_name}",
            )
            end_gap = value_label.left - (bar.left + bar.width)
            self.assertGreaterEqual(end_gap, 0)
            end_gaps.append(end_gap)

        self.assertEqual(len(set(end_gaps)), 1)

    def test_all_shapes_stay_inside_slide_and_visible_runs_use_jhenghei(self):
        deck = build_deck(
            rank_feature_importances(CURRENT_IMPORTANCES),
            build_insights(CURRENT_IMPORTANCES),
        )
        slide = deck.slides[0]

        for shape in slide.shapes:
            with self.subTest(shape=shape.name):
                self.assertGreaterEqual(shape.left, 0)
                self.assertGreaterEqual(shape.top, 0)
                self.assertLessEqual(shape.left + shape.width, deck.slide_width)
                self.assertLessEqual(shape.top + shape.height, deck.slide_height)

            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text:
                        continue
                    with self.subTest(shape=shape.name, text=run.text):
                        self.assertEqual(run.font.name, "Microsoft JhengHei")
                        east_asian = run._r.get_or_add_rPr().find(qn("a:ea"))
                        self.assertIsNotNone(east_asian)
                        self.assertEqual(
                            east_asian.get("typeface"),
                            "Microsoft JhengHei",
                        )

    def test_creation_pipeline_writes_a_valid_native_shape_powerpoint(self):
        result = create_feature_importance_powerpoint(
            self.model_path,
            self.output_path,
        )

        self.assertEqual(result, self.output_path)
        with zipfile.ZipFile(self.output_path) as package:
            self.assertIsNone(package.testzip())
            slide_xml = ElementTree.fromstring(
                package.read("ppt/slides/slide1.xml")
            )
        namespace = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
        }
        self.assertEqual(slide_xml.findall(".//a:blip", namespace), [])

        reopened = Presentation(self.output_path)
        self.assertEqual(len(reopened.slides), 1)
        self.assertAlmostEqual(
            reopened.slide_width / reopened.slide_height,
            16 / 9,
            places=2,
        )

    def test_save_failure_preserves_existing_output_and_cleans_temporary_file(self):
        self.output_path.write_bytes(b"existing deck")

        class FailingDeck:
            def save(self, _path):
                raise OSError("simulated save failure")

        with patch(
            "scripts.create_random_forest_feature_importance_powerpoint.build_deck",
            return_value=FailingDeck(),
        ):
            with self.assertRaisesRegex(OSError, "simulated save failure"):
                create_feature_importance_powerpoint(
                    self.model_path,
                    self.output_path,
                )

        self.assertEqual(self.output_path.read_bytes(), b"existing deck")
        self.assertEqual(
            list(self.root.glob(".feature-importance-*.pptx")),
            [],
        )

    def test_corrupt_package_preserves_existing_output_and_cleans_temporary_file(self):
        self.output_path.write_bytes(b"existing deck")

        class CorruptDeck:
            def save(self, path):
                Path(path).write_bytes(b"not a PowerPoint package")

        with patch(
            "scripts.create_random_forest_feature_importance_powerpoint.build_deck",
            return_value=CorruptDeck(),
        ):
            with self.assertRaisesRegex(ValueError, "corrupt"):
                create_feature_importance_powerpoint(
                    self.model_path,
                    self.output_path,
                )

        self.assertEqual(self.output_path.read_bytes(), b"existing deck")
        self.assertEqual(
            list(self.root.glob(".feature-importance-*.pptx")),
            [],
        )


if __name__ == "__main__":
    unittest.main()
