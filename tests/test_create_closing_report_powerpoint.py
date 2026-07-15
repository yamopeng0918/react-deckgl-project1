import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation

from scripts.create_closing_report_powerpoint import (
    DEFAULT_FRAME_NAMES,
    create_closing_report_deck,
    default_frame_paths,
)


class ClosingReportPowerPointTest(unittest.TestCase):
    def test_creates_fourteen_widescreen_slides_with_page_four_hyperlink(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            frames = []
            for index in range(14):
                frame = root / f"frame-{index:02d}.png"
                Image.new("RGB", (1920, 1080), (index * 10, 20, 40)).save(frame)
                frames.append(frame)
            linked_deck = root / "decision-tree-results.pptx"
            Presentation().save(linked_deck)
            output = root / "closing-report.pptx"

            create_closing_report_deck(frames, output, linked_deck.name)
            deck = Presentation(output)
            with zipfile.ZipFile(output) as package:
                slide_four_relationships = package.read(
                    "ppt/slides/_rels/slide4.xml.rels"
                ).decode("utf-8")

        self.assertEqual(len(deck.slides), 14)
        self.assertAlmostEqual(deck.slide_width / deck.slide_height, 16 / 9, places=2)
        self.assertIn("decision-tree-results.pptx", slide_four_relationships)
        overlay = deck.slides[3].shapes[-1]
        self.assertAlmostEqual(overlay.left.inches, 1.00, places=2)
        self.assertAlmostEqual(overlay.top.inches, 6.17, places=2)
        self.assertAlmostEqual(overlay.width.inches, 2.10, places=2)
        self.assertAlmostEqual(overlay.height.inches, 0.40, places=2)

    def test_uses_documented_midpoint_frame_order(self):
        root = Path("snapshots")
        self.assertEqual(
            default_frame_paths(root),
            [root / name for name in DEFAULT_FRAME_NAMES],
        )

    def test_rejects_invalid_linked_deck(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            frames = []
            for index in range(14):
                frame = root / f"frame-{index:02d}.png"
                Image.new("RGB", (1920, 1080)).save(frame)
                frames.append(frame)
            linked_deck = root / "decision-tree-results.pptx"
            linked_deck.write_bytes(b"not a pptx")

            with self.assertRaisesRegex(ValueError, "not a valid PPTX"):
                create_closing_report_deck(
                    frames, root / "closing-report.pptx", linked_deck.name
                )


if __name__ == "__main__":
    unittest.main()
