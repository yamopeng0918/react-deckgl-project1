import argparse
import csv
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(19, 35, 56)
INK = RGBColor(34, 43, 54)
MUTED = RGBColor(102, 114, 128)
PAPER = RGBColor(247, 245, 239)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(31, 128, 125)
AMBER = RGBColor(218, 147, 53)
PALE = RGBColor(226, 231, 230)
FONT = "Microsoft JhengHei"


def _set_text(shape, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _text_box(slide, x, y, w, h, text, size, **kwargs):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(shape, text, size, **kwargs)
    return shape


def _rounded_box(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _load_comparison(path):
    with Path(path).open(encoding="utf-8", newline="") as stream:
        return {row["model"]: row for row in csv.DictReader(stream)}


def _number(value):
    return None if value in (None, "") else float(value)


def _validate_comparison(metrics_by_model, comparison):
    for model, metrics in metrics_by_model.items():
        if model not in comparison:
            raise ValueError(f"comparison missing model {model}")
        row = comparison[model]
        recalls = [metrics["recall"][str(label)] for label in metrics["labels"]]
        expected = {
            "accuracy": metrics["accuracy"],
            "macro_recall": sum(value for value in recalls if value is not None)
            / sum(value is not None for value in recalls),
        }
        for label in metrics["labels"]:
            expected[f"recall_{label}"] = metrics["recall"][str(label)]
            expected[f"support_{label}"] = metrics["support"][str(label)]
        for field, value in expected.items():
            actual = _number(row.get(field))
            if value is None and actual is None:
                continue
            if actual is None or not math.isclose(actual, float(value), rel_tol=1e-9, abs_tol=1e-9):
                raise ValueError(f"comparison {model} {field} disagrees with metrics")


def _comparison_sentence(comparison):
    tree = comparison["decision_tree"]
    forest = comparison["random_forest"]
    tree_macro, forest_macro = float(tree["macro_recall"]), float(forest["macro_recall"])
    tree_accuracy, forest_accuracy = float(tree["accuracy"]), float(forest["accuracy"])
    macro_tied = math.isclose(tree_macro, forest_macro)
    accuracy_tied = math.isclose(tree_accuracy, forest_accuracy)
    if macro_tied and accuracy_tied:
        return "比較結論｜測試 Macro Recall 與 Accuracy 皆平手，兩者沒有明確領先。"
    if macro_tied:
        accuracy_leader = "隨機森林" if forest_accuracy > tree_accuracy else "決策樹"
        return f"比較結論｜測試 Macro Recall 平手；{accuracy_leader}在 Accuracy 領先。"
    macro_leader = "隨機森林" if forest_macro > tree_macro else "決策樹"
    if accuracy_tied:
        accuracy_comparison = "Accuracy 平手"
    else:
        accuracy_leader = "隨機森林" if forest_accuracy > tree_accuracy else "決策樹"
        accuracy_comparison = f"Accuracy 亦由{accuracy_leader}領先"
    return f"比較結論｜{macro_leader}在時間外推測試的 Macro Recall 領先；{accuracy_comparison}。"


def _model_card(slide, x, title, metrics, comparison, accent):
    _rounded_box(slide, x, 1.55, 3.75, 4.92, WHITE, PALE)
    _text_box(slide, x + 0.22, 1.73, 3.3, 0.31, title, 16, color=accent, bold=True)
    _text_box(slide, x + 0.22, 2.08, 1.62, 0.25, "Accuracy", 8.5, color=MUTED)
    _text_box(slide, x + 0.22, 2.31, 1.62, 0.36, f"{metrics['accuracy'] * 100:.2f}%", 18, bold=True)
    macro = float(comparison["macro_recall"])
    _text_box(slide, x + 1.93, 2.08, 1.60, 0.25, "Macro Recall", 8.5, color=MUTED)
    _text_box(slide, x + 1.93, 2.31, 1.60, 0.36, f"{macro * 100:.2f}%", 18, bold=True)
    _text_box(slide, x + 0.22, 2.83, 3.2, 0.25, "震度   Recall（support）", 9, color=NAVY, bold=True)
    for index, label in enumerate(metrics["labels"]):
        recall = metrics["recall"][str(label)]
        support = metrics["support"][str(label)]
        value = "N/A" if recall is None else f"{recall * 100:.1f}%"
        y = 3.13 + index * 0.31
        _text_box(slide, x + 0.24, y, 0.30, 0.21, str(label), 8.5, color=NAVY, bold=True)
        _text_box(slide, x + 0.62, y, 1.45, 0.21, f"{value}（{support:,}）", 8.5)
    params = metrics["selected_parameters"]
    selected = [f"{key}={params[key]}" for key in ("n_estimators", "max_depth", "min_samples_leaf") if key in params]
    _text_box(slide, x + 0.22, 5.72, 3.28, 0.45, "選定參數｜" + " · ".join(selected), 7.6, color=MUTED)
    if "validation_macro_recall" in params:
        _text_box(slide, x + 0.22, 6.14, 3.28, 0.20, f"驗證 Macro Recall {params['validation_macro_recall'] * 100:.2f}% · Accuracy {params['validation_accuracy'] * 100:.2f}%", 7.2, color=MUTED)


def create_results_slide(decision_tree_metrics_path, random_forest_metrics_path, comparison_path, output_path):
    decision_tree = json.loads(Path(decision_tree_metrics_path).read_text(encoding="utf-8"))
    random_forest = json.loads(Path(random_forest_metrics_path).read_text(encoding="utf-8"))
    comparison = _load_comparison(comparison_path)
    metrics_by_model = {"decision_tree": decision_tree, "random_forest": random_forest}
    _validate_comparison(metrics_by_model, comparison)
    for field in ("periods", "train_rows", "test_rows"):
        if decision_tree[field] != random_forest[field]:
            raise ValueError(f"models do not share common {field}")

    deck = Presentation()
    deck.slide_width, deck.slide_height = Inches(13.333333), Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    _rounded_box(slide, 0, 0, 13.333, 0.16, TEAL, TEAL, radius=False)
    _text_box(slide, 0.55, 0.34, 8.8, 0.48, "台灣地震最大震度｜兩模型時間外推比較", 24, color=NAVY, bold=True)
    train, test = decision_tree["periods"]["train"], decision_tree["periods"]["test"]
    _text_box(slide, 0.57, 0.87, 12.0, 0.25, f"共同評估｜訓練 {train[0]}–{train[1]}（{decision_tree['train_rows']:,} 筆）｜測試 {test[0]}–{test[1]}（{decision_tree['test_rows']:,} 筆）", 10.5, color=MUTED)
    _model_card(slide, 0.55, "決策樹", decision_tree, comparison["decision_tree"], TEAL)
    _model_card(slide, 4.45, "隨機森林", random_forest, comparison["random_forest"], AMBER)
    _rounded_box(slide, 8.35, 1.55, 4.43, 4.92, NAVY, NAVY)
    _text_box(slide, 8.66, 1.82, 3.8, 0.55, _comparison_sentence(comparison), 13, color=WHITE, bold=True)
    _text_box(slide, 8.66, 2.63, 3.78, 0.65, f"決策樹｜Macro Recall {float(comparison['decision_tree']['macro_recall']) * 100:.2f}%｜Accuracy {decision_tree['accuracy'] * 100:.2f}%", 10, color=WHITE)
    _text_box(slide, 8.66, 3.33, 3.78, 0.65, f"隨機森林｜Macro Recall {float(comparison['random_forest']['macro_recall']) * 100:.2f}%｜Accuracy {random_forest['accuracy'] * 100:.2f}%", 10, color=WHITE)
    _text_box(slide, 8.66, 4.35, 3.75, 1.35, "稀有類別限制｜測試集震度 0 僅 1 筆、震度 6 僅 4 筆、震度 7 為 0 筆；其 Recall 不宜視為穩定結論。", 11, color=WHITE, bold=True)
    _text_box(slide, 0.58, 6.79, 12.1, 0.31, "判讀原則｜先比較 chronological test Macro Recall，再以 Accuracy 作為次要排序。", 9, color=MUTED)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Create a one-slide two-model results deck")
    parser.add_argument("--decision-tree-metrics", type=Path, default=Path("data/model/decision_tree_metrics.json"))
    parser.add_argument("--random-forest-metrics", type=Path, default=Path("data/model/random_forest_metrics.json"))
    parser.add_argument("--comparison", type=Path, default=Path("data/model/model_comparison.csv"))
    parser.add_argument("--output", type=Path, default=Path("data/model/decision-tree-results-review.pptx"))
    arguments = parser.parse_args()
    print(create_results_slide(arguments.decision_tree_metrics, arguments.random_forest_metrics, arguments.comparison, arguments.output))


if __name__ == "__main__":
    main()
