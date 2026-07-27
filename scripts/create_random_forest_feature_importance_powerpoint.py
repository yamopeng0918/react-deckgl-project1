import argparse
import math
import os
import tempfile
import zipfile
from pathlib import Path

import joblib
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


FEATURE_NAMES = (
    "magnitude",
    "depth_km",
    "longitude",
    "latitude",
    "month",
    "hour",
)

FEATURE_LABELS = {
    "magnitude": "規模",
    "depth_km": "深度",
    "longitude": "經度",
    "latitude": "緯度",
    "month": "月份",
    "hour": "時刻",
}

NAVY = RGBColor(25, 42, 62)
INK = RGBColor(39, 52, 65)
MUTED = RGBColor(101, 113, 124)
PAPER = RGBColor(246, 247, 243)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(25, 133, 129)
CARD_LINE = RGBColor(220, 226, 224)
BAR_COLORS = (
    TEAL,
    RGBColor(63, 148, 144),
    RGBColor(92, 162, 158),
    RGBColor(120, 176, 172),
    RGBColor(148, 190, 186),
    RGBColor(175, 204, 200),
)
FONT = "Microsoft JhengHei"


def load_feature_importances(model_path) -> tuple[tuple[str, float], ...]:
    """Load and validate the fitted forest's six normalized importances."""
    model_path = Path(model_path)
    if not model_path.is_file():
        raise FileNotFoundError(f"model file not found: {model_path}")

    model = joblib.load(model_path)
    raw_importances = getattr(model, "feature_importances_", None)
    if raw_importances is None:
        raise ValueError("Model must expose feature_importances_")

    try:
        values = tuple(float(value) for value in raw_importances)
    except (TypeError, ValueError) as error:
        raise ValueError("feature importances must be finite non-negative numbers") from error

    if len(values) != len(FEATURE_NAMES):
        raise ValueError("Model must expose exactly six feature importances")
    if any(not math.isfinite(value) or value < 0 for value in values):
        raise ValueError("feature importances must be finite non-negative numbers")

    total = sum(values)
    if not math.isclose(total, 1.0, rel_tol=1e-6, abs_tol=1e-6):
        raise ValueError("feature importances must sum to one")

    return tuple(zip(FEATURE_NAMES, values))


def rank_feature_importances(importances) -> tuple[tuple[str, float], ...]:
    """Return importances from highest to lowest with stable feature-order ties."""
    feature_order = {name: index for index, name in enumerate(FEATURE_NAMES)}
    indexed_importances = tuple(importances)
    return tuple(
        sorted(
            indexed_importances,
            key=lambda item: (-item[1], feature_order.get(item[0], len(feature_order))),
        )
    )


def build_insights(importances) -> tuple[str, str, str]:
    """Build model-derived summary copy that remains accurate as values change."""
    values = dict(importances)
    top_name, top_value = rank_feature_importances(importances)[0]
    top_label = FEATURE_LABELS[top_name]
    spatial_total = values["longitude"] + values["latitude"]
    temporal_total = values["month"] + values["hour"]

    if spatial_total > top_value:
        spatial_comparison_text = (
            f"經緯度合計 {spatial_total:.2%}，顯示地理位置整體影響高於單一規模"
        )
    else:
        spatial_comparison_text = (
            f"經緯度合計 {spatial_total:.2%}，地理位置的相對影響應與其他特徵一併解讀"
        )

    magnitude_value = values["magnitude"]
    if temporal_total <= magnitude_value and temporal_total <= spatial_total:
        temporal_comparison_text = (
            f"月份與時刻合計 {temporal_total:.2%}，時間訊號存在，但不是主要依據"
        )
    else:
        temporal_comparison_text = (
            f"月份與時刻合計 {temporal_total:.2%}，時間訊號的相對影響應與其他特徵一併解讀"
        )

    return (
        f"{top_label}是最重要的單一特徵（{top_value:.2%}）",
        spatial_comparison_text,
        temporal_comparison_text,
    )


def _set_run_typeface(run):
    run.font.name = FONT
    properties = run._r.get_or_add_rPr()
    east_asian = properties.find(qn("a:ea"))
    if east_asian is None:
        east_asian = OxmlElement("a:ea")
        latin = properties.find(qn("a:latin"))
        if latin is None:
            properties.append(east_asian)
        else:
            latin.addnext(east_asian)
    east_asian.set("typeface", FONT)


def _add_text(
    slide,
    x,
    y,
    width,
    height,
    value,
    *,
    size,
    color,
    bold=False,
    alignment=PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = value
    _set_run_typeface(run)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_insight_card(slide, index, y, insight):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.15),
        Inches(y),
        Inches(3.63),
        Inches(1.12),
    )
    card.name = f"importance-insight-{index}"
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_LINE
    card.line.width = Pt(1)

    frame = card.text_frame
    frame.clear()
    frame.margin_left = Inches(0.26)
    frame.margin_right = Inches(0.16)
    frame.margin_top = Inches(0.10)
    frame.margin_bottom = Inches(0.10)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = insight
    _set_run_typeface(run)
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = INK

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(9.15),
        Inches(y),
        Inches(0.08),
        Inches(1.12),
    )
    accent.name = f"importance-insight-accent-{index}"
    accent.fill.solid()
    accent.fill.fore_color.rgb = BAR_COLORS[index - 1]
    accent.line.color.rgb = BAR_COLORS[index - 1]
    return card


def build_deck(importances, insights) -> Presentation:
    """Build a one-slide native-shape feature-importance chart."""
    importances = tuple(importances)
    insights = tuple(insights)
    if len(importances) != len(FEATURE_NAMES):
        raise ValueError("importances must contain exactly six features")
    if len(insights) != 3:
        raise ValueError("insights must contain exactly three strings")

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        deck.slide_width,
        Inches(0.12),
    )
    accent.name = "top-accent"
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.color.rgb = TEAL

    _add_text(
        slide,
        0.65,
        0.27,
        12.13,
        0.48,
        "什麼特徵最影響最大震度分類？",
        size=24,
        color=NAVY,
        bold=True,
    ).name = "slide-title"
    _add_text(
        slide,
        0.65,
        0.77,
        12.13,
        0.30,
        "隨機森林特徵重要度｜數值越高，模型越常依賴該特徵做判斷",
        size=10.5,
        color=MUTED,
    ).name = "slide-subtitle"

    maximum_importance = importances[0][1]
    if maximum_importance <= 0:
        raise ValueError("the strongest feature importance must be positive")

    baseline_x = 1.70
    maximum_width = 5.85
    row_positions = (1.55, 2.35, 3.15, 3.95, 4.75, 5.55)
    for index, ((feature_name, value), y) in enumerate(
        zip(importances, row_positions)
    ):
        label = _add_text(
            slide,
            0.65,
            y,
            0.85,
            0.36,
            FEATURE_LABELS[feature_name],
            size=11.5,
            color=INK,
            bold=index == 0,
        )
        label.name = f"importance-label-{feature_name}"

        bar_width = maximum_width * value / maximum_importance
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(baseline_x),
            Inches(y),
            Inches(bar_width),
            Inches(0.36),
        )
        bar.name = f"importance-bar-{feature_name}"
        bar.fill.solid()
        bar.fill.fore_color.rgb = BAR_COLORS[index]
        bar.line.color.rgb = BAR_COLORS[index]

        value_label = _add_text(
            slide,
            7.68,
            y,
            1.07,
            0.36,
            f"{value:.2%}",
            size=11,
            color=NAVY,
            bold=index == 0,
            alignment=PP_ALIGN.RIGHT,
        )
        value_label.name = f"importance-value-{feature_name}"

    _add_text(
        slide,
        9.15,
        1.11,
        3.63,
        0.25,
        "模型洞見",
        size=10,
        color=TEAL,
        bold=True,
    ).name = "insight-heading"
    for index, (y, insight) in enumerate(
        zip((1.48, 2.83, 4.18), insights),
        start=1,
    ):
        _add_insight_card(slide, index, y, insight)

    limitation = (
        "重要度代表模型使用程度，不等於因果關係；經緯度也可能共同反映區域差異。"
        "本模型是最大震度分類，不是地震預測。"
    )
    _add_text(
        slide,
        0.65,
        6.55,
        12.13,
        0.45,
        limitation,
        size=9.5,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    ).name = "limitation"
    return deck


def _verify_generated_deck(path):
    try:
        with zipfile.ZipFile(path) as package:
            if package.testzip() is not None:
                raise ValueError("generated PowerPoint package is corrupt")
        reopened = Presentation(path)
        if len(reopened.slides) != 1:
            raise ValueError("generated PowerPoint must contain exactly one slide")
        if not abs(reopened.slide_width / reopened.slide_height - 16 / 9) < 0.01:
            raise ValueError("generated PowerPoint must use a 16:9 slide size")
    except Exception as error:
        if isinstance(error, ValueError):
            raise
        raise ValueError(
            "generated PowerPoint package is corrupt or cannot be reopened"
        ) from error


def create_feature_importance_powerpoint(model_path, output_path) -> Path:
    """Load model importances and atomically write the one-slide chart."""
    importances = load_feature_importances(model_path)
    ranked = rank_feature_importances(importances)
    insights = build_insights(importances)
    deck = build_deck(ranked, insights)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.stem}-",
        suffix=".pptx",
        dir=output_path.parent,
    )
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        deck.save(temporary)
        _verify_generated_deck(temporary)
        os.replace(temporary, output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Create the random-forest feature-importance PowerPoint"
    )
    parser.add_argument(
        "--model",
        type=Path,
        default=Path("data/model/random_forest_model.joblib"),
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("data/model/random-forest-feature-importance.pptx"),
    )
    arguments = parser.parse_args()
    print(
        create_feature_importance_powerpoint(
            arguments.model,
            arguments.output,
        )
    )


if __name__ == "__main__":
    main()
