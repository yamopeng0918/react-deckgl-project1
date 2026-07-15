import argparse
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SLIDE_WIDTH = 13.333333
SLIDE_HEIGHT = 7.5

DEFAULT_FRAME_NAMES = [
    "frame-00-at-3.5s.png",
    "frame-01-at-10.5s.png",
    "frame-02-at-17.5s.png",
    "frame-03-at-24.5s.png",
    "frame-04-at-31.5s.png",
    "frame-05-at-39.0s.png",
    "frame-06-at-46.5s.png",
    "frame-07-at-54.0s.png",
    "frame-08-at-62.5s.png",
    "frame-09-at-70.5s.png",
    "frame-10-at-78.0s.png",
    "frame-11-at-84.5s.png",
    "frame-12-at-90.0s.png",
    "frame-13-at-94.5s.png",
]


def default_frame_paths(snapshot_dir):
    snapshot_dir = Path(snapshot_dir)
    return [snapshot_dir / name for name in DEFAULT_FRAME_NAMES]


def validate_linked_deck(linked_deck_path):
    linked_deck_path = Path(linked_deck_path)
    if not linked_deck_path.is_file():
        raise FileNotFoundError(f"Missing linked deck: {linked_deck_path}")
    if not zipfile.is_zipfile(linked_deck_path):
        raise ValueError(f"Linked deck is not a valid PPTX package: {linked_deck_path}")
    Presentation(linked_deck_path)


def create_closing_report_deck(frame_paths, output_path, linked_deck_name):
    frame_paths = [Path(path) for path in frame_paths]
    if len(frame_paths) != 14:
        raise ValueError(f"Expected 14 scene frames, got {len(frame_paths)}")
    missing = [str(path) for path in frame_paths if not path.is_file()]
    if missing:
        raise FileNotFoundError(f"Missing scene frames: {', '.join(missing)}")

    output_path = Path(output_path)
    linked_deck_path = Path(linked_deck_name)
    if not linked_deck_path.is_absolute():
        linked_deck_path = output_path.parent / linked_deck_path
    validate_linked_deck(linked_deck_path)

    deck = Presentation()
    deck.slide_width = Inches(SLIDE_WIDTH)
    deck.slide_height = Inches(SLIDE_HEIGHT)
    blank_layout = deck.slide_layouts[6]

    for index, frame_path in enumerate(frame_paths):
        slide = deck.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(frame_path),
            0,
            0,
            width=deck.slide_width,
            height=deck.slide_height,
        )

        if index == 3:
            link_overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.00),
                Inches(6.17),
                Inches(2.10),
                Inches(0.40),
            )
            link_overlay.fill.background()
            link_overlay.line.fill.background()
            link_overlay.click_action.hyperlink.address = Path(linked_deck_name).name

    output_path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Build a 14-slide PowerPoint from closing-report scene frames"
    )
    parser.add_argument("frames", nargs="*", type=Path)
    parser.add_argument(
        "--snapshot-dir",
        type=Path,
        default=Path("hyperframes/closing-report/snapshots"),
        help="Directory containing the 14 documented midpoint frames",
    )
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument(
        "--linked-deck", default="decision-tree-results.pptx"
    )
    arguments = parser.parse_args()
    if arguments.frames and len(arguments.frames) != 14:
        parser.error("provide either no frame paths or exactly 14 frame paths")
    frames = arguments.frames or default_frame_paths(arguments.snapshot_dir)
    output = create_closing_report_deck(frames, arguments.output, arguments.linked_deck)
    print(output)


if __name__ == "__main__":
    main()
