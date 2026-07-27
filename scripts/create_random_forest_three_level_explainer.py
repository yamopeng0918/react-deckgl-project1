import argparse
import os
import tempfile
import warnings
import zipfile
from collections import Counter
from pathlib import Path

import joblib
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.tree import DecisionTreeClassifier


FEATURE_NAMES = (
    "magnitude",
    "depth_km",
    "longitude",
    "latitude",
    "month",
    "hour",
)

FEATURE_LABELS = {
    "magnitude": "震級",
    "depth_km": "深度",
    "longitude": "經度",
    "latitude": "緯度",
    "month": "月份",
    "hour": "時刻",
}

LOCATION_FEATURE_INDICES = frozenset((2, 3))
EXPECTED_CLASS_LABELS = frozenset(range(8))

NAVY = RGBColor(19, 35, 56)
INK = RGBColor(34, 43, 54)
MUTED = RGBColor(102, 114, 128)
PAPER = RGBColor(247, 245, 239)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(31, 128, 125)
AMBER = RGBColor(218, 147, 53)
MIDDLE_BLUE = RGBColor(128, 154, 175)
PALE_BLUE = RGBColor(223, 231, 235)
FONT = "Microsoft JhengHei"

NODE_WIDTH = 2.95
NODE_HEIGHT = 0.82
NODE_POSITIONS = {
    0: ((5.19, 1.18),),
    1: ((2.15, 2.65), (8.23, 2.65)),
    2: ((0.45, 4.20), (3.45, 4.20), (6.45, 4.20), (9.45, 4.20)),
}


def _validate_model(model):
    estimators = getattr(model, "estimators_", None)
    if not estimators:
        raise ValueError("Model must expose non-empty estimators_")

    if getattr(model, "n_features_in_", None) != len(FEATURE_NAMES):
        raise ValueError("Model must have six input features")

    classes = getattr(model, "classes_", None)
    if classes is None or len(classes) == 0:
        raise ValueError("Model must expose non-empty classes_")
    if not set(classes).issubset(EXPECTED_CLASS_LABELS):
        raise ValueError("Model classes must be compatible with 0-7")

    return estimators


def _nodes_through_depth(tree, max_depth=2):
    nodes = []
    pending = [(0, 0)]
    tree_data = tree.tree_
    while pending:
        node_id, depth = pending.pop()
        nodes.append((node_id, depth))
        if depth == max_depth:
            continue
        left_child = int(tree_data.children_left[node_id])
        if left_child == -1:
            continue
        right_child = int(tree_data.children_right[node_id])
        pending.append((right_child, depth + 1))
        pending.append((left_child, depth + 1))
    return nodes


def _count_nodes_through_depth_two(tree):
    return len(_nodes_through_depth(tree))


def _count_location_splits_through_depth_two(tree):
    tree_data = tree.tree_
    return sum(
        int(tree_data.feature[node_id]) in LOCATION_FEATURE_INDICES
        for node_id, _ in _nodes_through_depth(tree)
        if int(tree_data.children_left[node_id]) != -1
    )


def select_representative_tree(model) -> tuple[int, DecisionTreeClassifier]:
    """Return the most readable deterministic representative forest tree."""
    estimators = _validate_model(model)
    root_features = [int(estimator.tree_.feature[0]) for estimator in estimators]
    most_common_root_feature = Counter(root_features).most_common(1)[0][0]

    ranked_estimators = []
    for estimator_index, estimator in enumerate(estimators):
        root_feature = int(estimator.tree_.feature[0])
        rank = (
            root_feature != most_common_root_feature,
            -_count_nodes_through_depth_two(estimator),
            _count_location_splits_through_depth_two(estimator),
            estimator_index,
        )
        ranked_estimators.append((rank, estimator_index, estimator))

    _, estimator_index, estimator = min(ranked_estimators)
    return estimator_index, estimator


def extract_three_levels(estimator, feature_names, class_labels) -> list[dict]:
    """Extract root through depth two in left-to-right traversal order."""
    if not class_labels:
        raise ValueError("class_labels must not be empty")

    tree_data = estimator.tree_
    nodes = []
    pending = [(0, None, None, 0)]
    while pending:
        node_id, parent_id, branch, depth = pending.pop()
        left_child = int(tree_data.children_left[node_id])
        is_leaf = left_child == -1
        feature_index = None if is_leaf else int(tree_data.feature[node_id])
        class_counts = tuple(float(value) for value in tree_data.value[node_id][0])
        class_index = max(range(len(class_counts)), key=class_counts.__getitem__)

        nodes.append(
            {
                "node_id": int(node_id),
                "parent_id": parent_id,
                "branch": branch,
                "depth": depth,
                "feature_index": feature_index,
                "feature_name": (
                    None if feature_index is None else feature_names[feature_index]
                ),
                "threshold": (
                    None
                    if feature_index is None
                    else float(tree_data.threshold[node_id])
                ),
                "samples": int(tree_data.n_node_samples[node_id]),
                "dominant_class": int(class_labels[class_index]),
                "class_counts": class_counts,
                "is_leaf": is_leaf,
            }
        )

        if is_leaf or depth == 2:
            continue
        right_child = int(tree_data.children_right[node_id])
        pending.append((right_child, node_id, "no", depth + 1))
        pending.append((left_child, node_id, "yes", depth + 1))

    return nodes


def _add_text(
    slide,
    x,
    y,
    width,
    height,
    value,
    *,
    size,
    color,
    bold=False,
    alignment=PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_box(slide, x, y, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _node_copy(node):
    if node["is_leaf"]:
        return (
            f"此分支預測震度 {node['dominant_class']}\n"
            f"樣本 {node['samples']:,}"
        )
    return (
        f"{FEATURE_LABELS[node['feature_name']]} ≤ {node['threshold']:.2f}？\n"
        f"樣本 {node['samples']:,}｜目前偏向震度 "
        f"{node['dominant_class']}"
    )


def _node_slots(nodes):
    slots = {}
    by_id = {node["node_id"]: node for node in nodes}
    for node in nodes:
        if node["depth"] == 0:
            slots[node["node_id"]] = 0
            continue
        parent = by_id[node["parent_id"]]
        parent_slot = slots[parent["node_id"]]
        slots[node["node_id"]] = parent_slot * 2 + (node["branch"] == "no")
    return slots


def _add_branch_label(slide, parent_position, child_position, branch):
    parent_x, parent_y = parent_position
    child_x, child_y = child_position
    color = TEAL if branch == "yes" else AMBER
    label = "是（≤）" if branch == "yes" else "否（>）"
    center_x = ((parent_x + NODE_WIDTH / 2) + (child_x + NODE_WIDTH / 2)) / 2
    center_y = ((parent_y + NODE_HEIGHT) + child_y) / 2
    return _add_text(
        slide,
        center_x - 0.37,
        center_y - 0.12,
        0.74,
        0.24,
        label,
        size=8.5,
        color=color,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def build_deck(nodes, tree_index, estimator_count) -> Presentation:
    """Build one native-shape slide explaining a representative tree."""
    if not nodes or nodes[0]["depth"] != 0:
        raise ValueError("nodes must begin with a root node")
    if estimator_count <= 0:
        raise ValueError("estimator_count must be positive")

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        deck.slide_width,
        Inches(0.12),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.color.rgb = TEAL

    _add_text(
        slide,
        0.55,
        0.25,
        8.7,
        0.48,
        "隨機森林如何判斷最大震度？",
        size=24,
        color=NAVY,
        bold=True,
    )
    _add_text(
        slide,
        0.57,
        0.73,
        11.8,
        0.28,
        f"代表樹第 {tree_index + 1} 棵｜以下為模型真實前三層判斷規則",
        size=10,
        color=MUTED,
    )

    slots = _node_slots(nodes)
    node_positions = {
        node["node_id"]: NODE_POSITIONS[node["depth"]][slots[node["node_id"]]]
        for node in nodes
    }

    for node in nodes:
        if node["parent_id"] is None:
            continue
        parent_x, parent_y = node_positions[node["parent_id"]]
        child_x, child_y = node_positions[node["node_id"]]
        color = TEAL if node["branch"] == "yes" else AMBER
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(parent_x + NODE_WIDTH / 2),
            Inches(parent_y + NODE_HEIGHT),
            Inches(child_x + NODE_WIDTH / 2),
            Inches(child_y),
        )
        connector.line.color.rgb = color
        connector.line.width = Pt(2)
        _add_branch_label(
            slide,
            (parent_x, parent_y),
            (child_x, child_y),
            node["branch"],
        )

    fills = {0: NAVY, 1: MIDDLE_BLUE, 2: PALE_BLUE}
    text_colors = {0: WHITE, 1: WHITE, 2: INK}
    for node in nodes:
        x, y = node_positions[node["node_id"]]
        box = _add_box(
            slide,
            x,
            y,
            NODE_WIDTH,
            NODE_HEIGHT,
            fills[node["depth"]],
            NAVY if node["depth"] == 2 else fills[node["depth"]],
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(0.12)
        frame.margin_right = Inches(0.12)
        frame.margin_top = Inches(0.05)
        frame.margin_bottom = Inches(0.05)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = _node_copy(node)
        run.font.name = FONT
        run.font.size = Pt(10.5 if node["depth"] == 2 else 11)
        run.font.bold = True
        run.font.color.rgb = text_colors[node["depth"]]

    _add_box(slide, 0.55, 5.37, 12.23, 0.92, NAVY)
    _add_text(
        slide,
        0.85,
        5.51,
        11.63,
        0.62,
        (
            f"這只是 {estimator_count} 棵樹中的一棵。每棵樹各自判斷，"
            "最後由所有樹投票決定分類結果。"
        ),
        size=13,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        0.65,
        6.55,
        12.03,
        0.36,
        (
            "前三層只呈現模型早期的主要切分，不代表完整模型；"
            "本模型是最大震度分類，不是地震預測。"
        ),
        size=9.5,
        color=MUTED,
        alignment=PP_ALIGN.CENTER,
    )
    return deck


def _verify_generated_deck(path):
    try:
        with zipfile.ZipFile(path) as package:
            if package.testzip() is not None:
                raise ValueError("generated PowerPoint package is corrupt")
        reopened = Presentation(path)
        if len(reopened.slides) != 1:
            raise ValueError("generated PowerPoint must contain exactly one slide")
        if not abs(reopened.slide_width / reopened.slide_height - 16 / 9) < 0.01:
            raise ValueError("generated PowerPoint must use a 16:9 slide size")
    except Exception as error:
        if isinstance(error, ValueError):
            raise
        raise ValueError(
            "generated PowerPoint package is corrupt or cannot be reopened"
        ) from error


def create_random_forest_three_level_explainer(model_path, output_path) -> Path:
    """Load the fitted forest and atomically write its one-slide explainer."""
    with warnings.catch_warnings():
        warnings.filterwarnings(
            "ignore",
            message="Setting the shape on a NumPy array has been deprecated.*",
            category=DeprecationWarning,
        )
        model = joblib.load(model_path)
    tree_index, estimator = select_representative_tree(model)
    nodes = extract_three_levels(estimator, FEATURE_NAMES, list(model.classes_))
    deck = build_deck(nodes, tree_index, len(model.estimators_))

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix=f".{output_path.stem}-",
        suffix=".pptx",
        dir=output_path.parent,
    )
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        deck.save(temporary)
        _verify_generated_deck(temporary)
        os.replace(temporary, output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="Create the one-slide random-forest three-level explainer"
    )
    parser.add_argument(
        "--model",
        type=Path,
        default=Path("data/model/random_forest_model.joblib"),
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("data/model/random-forest-three-level-explainer.pptx"),
    )
    arguments = parser.parse_args()
    print(
        create_random_forest_three_level_explainer(
            arguments.model,
            arguments.output,
        )
    )


if __name__ == "__main__":
    main()
