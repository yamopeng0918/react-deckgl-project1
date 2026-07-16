import argparse
import io
import json
import math
import os
import tempfile
import zipfile
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", str(Path(tempfile.gettempdir()) / "taiwan-earthquake-matplotlib"))
import matplotlib
matplotlib.use("Agg")
from matplotlib import pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(19, 35, 56)
INK = RGBColor(34, 43, 54)
MUTED = RGBColor(102, 114, 128)
PAPER = RGBColor(247, 245, 239)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(31, 128, 125)
AMBER = RGBColor(218, 147, 53)
FONT = "Microsoft JhengHei"
PARAMETER_KEYS = {
    "n_estimators", "max_depth", "min_samples_leaf", "max_features",
    "class_weight", "random_state", "validation_start_year",
    "validation_end_year", "validation_macro_recall", "validation_accuracy",
}


def _is_finite_number(value):
    return not isinstance(value, bool) and isinstance(value, (int, float)) and math.isfinite(value)


def _is_integer(value):
    return not isinstance(value, bool) and isinstance(value, int)


def load_and_validate_metrics(metrics_path):
    metrics_path = Path(metrics_path)
    if not metrics_path.is_file():
        raise FileNotFoundError(f"metrics JSON not found: {metrics_path}")
    metrics = json.loads(metrics_path.read_text(encoding="utf-8"))
    labels = list(range(8))
    if metrics.get("labels") != labels:
        raise ValueError("labels must be exactly 0 through 7")
    expected_keys = {str(label) for label in labels}
    for field in ("recall", "support"):
        values = metrics.get(field)
        if not isinstance(values, dict) or set(values) != expected_keys:
            raise ValueError(f"{field} keys must be exactly 0 through 7")
    for label, value in metrics["recall"].items():
        if value is not None and (not _is_finite_number(value) or not 0 <= value <= 1):
            raise ValueError(f"recall {label} must be None or a finite number from 0 through 1")
    for label, value in metrics["support"].items():
        if not _is_integer(value) or value < 0:
            raise ValueError(f"support {label} must be a finite non-negative integer")
    non_null_recalls = 0
    for label in expected_keys:
        recall = metrics["recall"][label]
        support = metrics["support"][label]
        if support > 0 and recall is None:
            raise ValueError(f"recall {label} must be numeric when support is positive")
        if support == 0 and recall is not None:
            raise ValueError(f"recall {label} must be None when support is zero")
        non_null_recalls += recall is not None
    if non_null_recalls == 0:
        raise ValueError("at least one non-null recall is required")
    matrix = metrics.get("confusion_matrix")
    if not isinstance(matrix, list) or len(matrix) != 8 or any(
        not isinstance(row, list) or len(row) != 8 for row in matrix
    ):
        raise ValueError("confusion matrix must be exactly 8 x 8")
    for row_index, row in enumerate(matrix):
        if any(not _is_integer(value) or value < 0 for value in row):
            raise ValueError("confusion matrix must contain finite non-negative integer counts")
        support = metrics["support"][str(row_index)]
        if sum(row) != support:
            raise ValueError(f"confusion matrix row {row_index} total disagrees with support")
        if support > 0:
            expected_recall = row[row_index] / support
            if not math.isclose(metrics["recall"][str(row_index)], expected_recall, rel_tol=1e-12, abs_tol=1e-12):
                raise ValueError(f"recall {row_index} disagrees with confusion matrix diagonal/support")
    accuracy = metrics.get("accuracy")
    if not _is_finite_number(accuracy) or not 0 <= accuracy <= 1:
        raise ValueError("accuracy must be a finite number from 0 through 1")
    total_support = sum(metrics["support"].values())
    expected_accuracy = sum(matrix[label][label] for label in labels) / total_support
    if not math.isclose(accuracy, expected_accuracy, rel_tol=1e-12, abs_tol=1e-12):
        raise ValueError("accuracy disagrees with confusion matrix diagonal/total support")
    periods = metrics.get("periods")
    if not isinstance(periods, dict):
        raise ValueError("periods must contain ordered train and test year pairs")
    for name in ("train", "test"):
        period = periods.get(name)
        if (
            not isinstance(period, list) or len(period) != 2
            or any(not _is_integer(year) for year in period)
            or period[0] > period[1]
        ):
            raise ValueError(f"periods {name} must be two ordered integer years")
    for field in ("train_rows", "test_rows"):
        value = metrics.get(field)
        if not _is_integer(value) or value <= 0:
            raise ValueError(f"{field} must be a positive integer")
    parameters = metrics.get("selected_parameters")
    if not isinstance(parameters, dict) or not PARAMETER_KEYS.issubset(parameters):
        missing = sorted(PARAMETER_KEYS - set(parameters)) if isinstance(parameters, dict) else sorted(PARAMETER_KEYS)
        raise ValueError(f"selected_parameters must be a dict containing required keys; missing: {missing}")
    if not _is_integer(parameters["n_estimators"]) or parameters["n_estimators"] <= 0:
        raise ValueError("n_estimators must be a positive integer")
    max_depth = parameters["max_depth"]
    if max_depth is not None and (not _is_integer(max_depth) or max_depth <= 0):
        raise ValueError("max_depth must be None or a positive integer")
    if not _is_integer(parameters["min_samples_leaf"]) or parameters["min_samples_leaf"] <= 0:
        raise ValueError("min_samples_leaf must be a positive integer")
    if parameters["max_features"] != "sqrt":
        raise ValueError("max_features must be exactly 'sqrt'")
    if parameters["class_weight"] != "balanced_subsample":
        raise ValueError("class_weight must be exactly 'balanced_subsample'")
    if not _is_integer(parameters["random_state"]) or parameters["random_state"] != 42:
        raise ValueError("random_state must be the integer 42")
    start, end = parameters["validation_start_year"], parameters["validation_end_year"]
    if not _is_integer(start):
        raise ValueError("validation_start_year must be an integer")
    if not _is_integer(end):
        raise ValueError("validation_end_year must be an integer")
    if start > end:
        raise ValueError("validation_start_year must not exceed validation_end_year")
    for field in ("validation_macro_recall", "validation_accuracy"):
        value = parameters[field]
        if not _is_finite_number(value) or not 0 <= value <= 1:
            raise ValueError(f"{field} must be a finite number from 0 through 1")
    return metrics


def _render_confusion_matrix_png(metrics):
    matrix = metrics["confusion_matrix"]
    figure, axis = plt.subplots(figsize=(8, 6), dpi=150)
    image = axis.imshow(matrix, cmap="Blues", interpolation="nearest")
    axis.set_xlabel("Predicted intensity")
    axis.set_ylabel("Actual intensity")
    axis.set_xticks(range(8), labels=range(8))
    axis.set_yticks(range(8), labels=range(8))
    threshold = max(max(row) for row in matrix) / 2
    for row in range(8):
        for column in range(8):
            axis.text(column, row, str(matrix[row][column]), ha="center", va="center",
                      color="white" if matrix[row][column] > threshold else "#132338", fontsize=8)
    figure.colorbar(image, ax=axis, fraction=0.046, pad=0.04)
    figure.tight_layout()
    stream = io.BytesIO()
    figure.savefig(stream, format="png", facecolor="#F7F5EF", metadata={"Software": "taiwan-earthquake-visualization"})
    plt.close(figure)
    return stream.getvalue()


def _text(slide, x, y, w, h, value, size=12, color=INK, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = frame.paragraphs[0].add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _box(slide, x, y, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def _background(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.color.rgb = TEAL


def _add_contained_picture(slide, png_bytes, x, y, w, h):
    with Image.open(io.BytesIO(png_bytes)) as image:
        ratio = image.width / image.height
    frame_ratio = w / h
    if ratio > frame_ratio:
        width, height = w, w / ratio
    else:
        width, height = h * ratio, h
    return slide.shapes.add_picture(
        io.BytesIO(png_bytes), Inches(x + (w - width) / 2), Inches(y + (h - height) / 2),
        width=Inches(width), height=Inches(height),
    )


def _build_deck(metrics, _matrix_path=None):
    deck = Presentation()
    deck.slide_width, deck.slide_height = Inches(13.333333), Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _background(slide)
    _text(slide, 0.55, 0.32, 9.5, 0.5, "隨機森林最大震度分類結果", 25, NAVY, True)
    _text(slide, 0.57, 0.82, 12, 0.3, "依時間切分的獨立測試結果｜分類，不是地震預測", 10, MUTED)
    _box(slide, 0.55, 1.35, 3.7, 1.35, WHITE)
    _text(slide, 0.82, 1.57, 1.4, 0.28, "Accuracy", 11, MUTED, True)
    _text(slide, 0.82, 1.9, 2.6, 0.48, f"{metrics['accuracy'] * 100:.2f}%", 26, TEAL, True)
    recalls = [value for value in metrics["recall"].values() if value is not None]
    macro = sum(recalls) / len(recalls)
    _box(slide, 4.45, 1.35, 3.7, 1.35, WHITE)
    _text(slide, 4.72, 1.57, 2.0, 0.28, "Macro Recall", 11, MUTED, True)
    _text(slide, 4.72, 1.9, 2.6, 0.48, f"{macro * 100:.2f}%", 26, AMBER, True)
    train, test = metrics["periods"]["train"], metrics["periods"]["test"]
    _text(slide, 8.55, 1.45, 4.1, 0.5, f"訓練 {train[0]}–{train[1]}｜{metrics['train_rows']:,} 筆", 12, NAVY, True)
    _text(slide, 8.55, 2.02, 4.1, 0.5, f"測試 {test[0]}–{test[1]}｜{metrics['test_rows']:,} 筆", 12, NAVY, True)
    _text(slide, 0.58, 2.97, 4.1, 0.34, "各震度 Recall / support", 14, NAVY, True)
    for label in range(8):
        recall, support = metrics["recall"][str(label)], metrics["support"][str(label)]
        value = "N/A" if support == 0 else f"{recall * 100:.1f}%"
        column, row = divmod(label, 4)
        _text(slide, 0.62 + column * 3.1, 3.4 + row * 0.58, 2.85, 0.42,
              f"震度 {label}：{value}（support {support:,}）", 10.5, INK, label in (0, 6, 7))
    params = metrics["selected_parameters"]
    parameter_text = "｜".join(f"{key}={value}" for key, value in params.items())
    _text(slide, 6.8, 3.0, 5.9, 1.55, "選定參數\n" + parameter_text, 9.2, NAVY, True)
    _box(slide, 6.8, 4.8, 5.95, 1.45, NAVY)
    _text(slide, 7.08, 4.98, 5.4, 1.02, "稀有類別警告：測試集震度 0 僅 1 筆、震度 6 僅 4 筆、震度 7 為 0 筆；相關 Recall 不宜視為穩定結論。", 11, WHITE, True)

    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _background(slide)
    _text(slide, 0.55, 0.3, 8.5, 0.48, "隨機森林混淆矩陣", 24, NAVY, True)
    _text(slide, 0.58, 0.78, 8.7, 0.3, "實際類別（縱軸）｜預測類別（橫軸）", 11, MUTED)
    _add_contained_picture(slide, _render_confusion_matrix_png(metrics), 0.55, 1.18, 8.35, 5.85)
    _box(slide, 9.25, 1.25, 3.55, 2.1, WHITE)
    _text(slide, 9.55, 1.5, 2.95, 1.55, "主要觀察\n2–4 級相鄰混淆最明顯，顯示模型容易在相近震度間錯分。", 13, NAVY, True)
    _box(slide, 9.25, 3.7, 3.55, 2.2, NAVY)
    _text(slide, 9.55, 3.98, 2.95, 1.55, "支持數警告\n0、6、7 級樣本極少或不存在，矩陣與 Recall 均須審慎解讀。", 12, WHITE, True)
    return deck


def create_random_forest_results_powerpoint(metrics_path, matrix_path, output_path):
    metrics = load_and_validate_metrics(metrics_path)
    matrix_path, output_path = Path(matrix_path), Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    deck = _build_deck(metrics, matrix_path)
    handle, temporary_name = tempfile.mkstemp(prefix=f".{output_path.stem}-", suffix=".pptx", dir=output_path.parent)
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        deck.save(temporary)
        try:
            with zipfile.ZipFile(temporary) as package:
                if package.testzip() is not None:
                    raise ValueError("generated PowerPoint package is corrupt")
            reopened = Presentation(temporary)
            if len(reopened.slides) != 2:
                raise ValueError("generated PowerPoint must contain exactly two slides")
        except Exception as error:
            raise ValueError("generated PowerPoint package is corrupt or cannot be reopened") from error
        os.replace(temporary, output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Create the two-slide random-forest results PowerPoint")
    parser.add_argument("--metrics", type=Path, default=Path("data/model/random_forest_metrics.json"))
    parser.add_argument("--matrix", type=Path, default=Path("data/model/random_forest_confusion_matrix.png"),
                        help="Deprecated compatibility option; ignored because the matrix is rendered from metrics JSON")
    parser.add_argument("--output", type=Path, default=Path("data/model/random-forest-results.pptx"))
    arguments = parser.parse_args()
    print(create_random_forest_results_powerpoint(arguments.metrics, arguments.matrix, arguments.output))


if __name__ == "__main__":
    main()
